"""
Generate hackathon submission files:
1. Detailed Document (.docx)
2. Presentation (.pptx)
3. Code archive (.zip)
"""
import os
import zipfile
from pathlib import Path

# ============================================================
# 1. DETAILED DOCUMENT (.docx)
# ============================================================
def create_detailed_document():
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    doc = Document()

    # -- Styles --
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)

    # -- Title Page --
    for _ in range(6):
        doc.add_paragraph()

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('ROADSoS AI')
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(220, 38, 38)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run('AI-Powered Emergency Response & Disaster Assistance Platform')
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(100, 100, 100)

    doc.add_paragraph()
    tagline = doc.add_paragraph()
    tagline.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = tagline.add_run('Road Safety Hackathon 2026 - BIMSTEC International Track')
    run.font.size = Pt(13)
    run.font.italic = True

    doc.add_paragraph()
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = meta.add_run('Submitted via Unstop | May 2026')
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(120, 120, 120)

    doc.add_page_break()

    # -- Table of Contents --
    doc.add_heading('Table of Contents', level=1)
    toc_items = [
        '1. Executive Summary',
        '2. Problem Statement',
        '3. Our Solution - ROADSoS AI',
        '4. Key Features',
        '5. System Architecture',
        '6. Technology Stack',
        '7. AI Integration - Gemini API',
        '8. Offline-First PWA Strategy',
        '9. User Workflows',
        '10. Impact & Scalability',
        '11. Future Roadmap',
        '12. Team & Acknowledgments',
    ]
    for item in toc_items:
        p = doc.add_paragraph(item)
        p.paragraph_format.space_after = Pt(2)
    doc.add_page_break()

    # -- 1. Executive Summary --
    doc.add_heading('1. Executive Summary', level=1)
    doc.add_paragraph(
        'ROADSoS AI is an AI-powered Progressive Web Application (PWA) designed to save lives during '
        'the "Golden Hour" - the critical first 60 minutes after a road accident or natural disaster. '
        'The platform combines real-time geolocation, Google Gemini AI for severity assessment and '
        'first-aid guidance, and offline-first architecture to ensure emergency assistance is available '
        'even in areas with poor or no internet connectivity.'
    )
    doc.add_paragraph(
        'According to the WHO, approximately 1.19 million people die each year from road traffic crashes, '
        'with over 90% of these deaths occurring in low- and middle-income countries. The BIMSTEC region '
        '(Bangladesh, Bhutan, India, Myanmar, Nepal, Sri Lanka, Thailand) accounts for a significant '
        'proportion of these fatalities. Many of these deaths are preventable if victims receive timely '
        'medical attention within the Golden Hour.'
    )
    doc.add_paragraph(
        'ROADSoS AI addresses this challenge by providing: (a) one-tap SOS activation with automatic GPS '
        'lock, (b) AI-powered severity assessment and step-by-step first-aid instructions with audio '
        'guidance, (c) real-time mapping of the nearest hospitals, police stations, and fire stations, '
        '(d) automatic crash detection using device accelerometer data, (e) offline emergency packs that '
        'cache critical medical data locally, and (f) shareable family tracking links for real-time '
        'incident monitoring.'
    )

    # -- 2. Problem Statement --
    doc.add_heading('2. Problem Statement', level=1)
    doc.add_paragraph(
        'Road accidents remain one of the leading causes of death and disability across the BIMSTEC nations. '
        'The critical challenges include:'
    )
    problems = [
        ('Delayed Response Times', 'In rural and semi-urban areas, emergency response times often exceed 30-45 minutes, well past the Golden Hour window.'),
        ('Lack of First-Aid Knowledge', 'Bystanders at accident scenes frequently lack basic first-aid skills, leading to improper handling that worsens injuries.'),
        ('Poor Connectivity', 'Many accident-prone highway stretches and rural roads have limited or no cellular connectivity, making it impossible to call for help.'),
        ('Information Gaps', 'First responders often arrive without critical medical information about the victim (blood type, allergies, existing conditions).'),
        ('Resource Discovery', 'Victims and bystanders struggle to identify the nearest medical facilities, especially in unfamiliar areas.'),
        ('Family Communication', 'Families are left in the dark during emergencies, with no real-time visibility into the victim\'s status or location.'),
    ]
    for title_text, desc in problems:
        p = doc.add_paragraph()
        run = p.add_run(f'{title_text}: ')
        run.bold = True
        p.add_run(desc)

    # -- 3. Our Solution --
    doc.add_heading('3. Our Solution - ROADSoS AI', level=1)
    doc.add_paragraph(
        'ROADSoS AI is a comprehensive, end-to-end emergency response platform that transforms any '
        'smartphone into a life-saving device. The solution operates across three critical phases:'
    )

    doc.add_heading('Phase 1: Detection & Activation', level=2)
    doc.add_paragraph(
        'The platform offers two activation methods: (a) Manual one-tap SOS with a 5-second countdown '
        'and GPS lock, allowing the user to confirm or cancel, and (b) Automatic crash detection using '
        'the device\'s accelerometer (DeviceMotionEvent API) to detect sudden deceleration patterns '
        'consistent with vehicular accidents. Upon activation, the system captures GPS coordinates, '
        'timestamps, speed, heading, network status, and battery level as a "Black Box" record.'
    )

    doc.add_heading('Phase 2: AI Assessment & Resource Discovery', level=2)
    doc.add_paragraph(
        'Once an SOS is triggered, Google Gemini AI (gemini-2.0-flash) analyzes the incident context '
        'and provides: a severity rating on a 1-5 scale, categorized first-aid instructions specific '
        'to the situation, audio narration of first-aid steps via Web Speech API, and a structured '
        'incident summary for emergency services. Simultaneously, the Overpass API queries '
        'OpenStreetMap data to locate hospitals, police stations, and fire stations within a '
        'configurable radius (default 5 km).'
    )

    doc.add_heading('Phase 3: Communication & Tracking', level=2)
    doc.add_paragraph(
        'The platform generates a shareable tracking link (/track/<incident_id>) that can be sent to '
        'family members via SMS, WhatsApp, or any messaging app. This public page requires no login '
        'and shows the victim\'s real-time location on a Leaflet map, along with AI-generated response '
        'instructions and incident status updates.'
    )

    # -- 4. Key Features --
    doc.add_heading('4. Key Features', level=1)

    features = [
        ('One-Tap SOS Activation', 'Large, accessible SOS button with 5-second countdown timer, GPS lock indicator, and haptic feedback. Designed for use under stress and panic.'),
        ('AI Emergency Assistant', 'Powered by Google Gemini API - provides real-time severity assessment (1-5 scale), categorized first-aid instructions, and audio guidance via text-to-speech.'),
        ('Auto Crash Detection', 'Uses the DeviceMotionEvent API to monitor accelerometer data. Detects sudden G-force spikes (>4G threshold) and automatically triggers SOS after a confirmation window.'),
        ('Emergency Resource Finder', 'Integrates with the Overpass API (OpenStreetMap) to locate nearby hospitals, police stations, and fire stations. Results displayed as color-coded markers on an interactive Leaflet.js map.'),
        ('Offline Emergency Pack', 'PWA architecture with Service Worker caching of UI assets and localStorage caching of medical profile and nearby resources. The app remains functional even without internet.'),
        ('Medical Profile & QR Card', 'Users can store blood group, allergies, medications, medical conditions, and organ donor status. A QR code is generated containing this data for quick scanning by first responders.'),
        ('Family Tracking Link', 'Public, no-login-required tracking page showing victim location, AI instructions, and incident status. Shareable via any messaging platform.'),
        ('Black Box Recorder', 'Automatically captures incident metadata including GPS coordinates, speed, heading, network status, battery level, and timestamps.'),
        ('Disaster Mode', 'Specialized interface for natural disasters (flood, earthquake, cyclone, landslide, fire, industrial accident) with tailored AI responses and shelter finding.'),
        ('Admin Dashboard & Heatmap', 'Administrative panel with user management, incident oversight, and a Leaflet.heat visualization showing incident density patterns for urban planning insights.'),
        ('Report Generation', 'Generates comprehensive incident reports in PDF, TXT, and JSON formats using ReportLab, suitable for insurance claims and legal documentation.'),
        ('100-Point Readiness Score', 'Gamified readiness scoring system encouraging users to complete their medical profile, add emergency contacts, generate QR cards, and download offline packs.'),
    ]
    for feat_title, feat_desc in features:
        p = doc.add_paragraph()
        run = p.add_run(f'{feat_title}: ')
        run.bold = True
        p.add_run(feat_desc)

    # -- 5. System Architecture --
    doc.add_heading('5. System Architecture', level=1)
    doc.add_paragraph(
        'ROADSoS AI follows a layered architecture pattern with clear separation of concerns:'
    )
    arch_layers = [
        ('Presentation Layer', 'HTML5 templates rendered via Jinja2, styled with TailwindCSS, with Leaflet.js for interactive mapping and vanilla JavaScript for SOS logic, crash detection, and PWA functionality.'),
        ('API Layer', 'FastAPI REST endpoints handling SOS activation, resource discovery, user management, AI interactions, admin operations, and report generation. JWT-based authentication with bcrypt password hashing.'),
        ('Service Layer', 'Business logic modules for user management (readiness scoring, password hashing), incident lifecycle management (creation, updates, resolution), and report generation (PDF/TXT/JSON via ReportLab).'),
        ('Data Layer', 'SQLAlchemy ORM with SQLite database. Models include User, MedicalProfile, EmergencyContact, Incident, BlackBoxRecord, CommunityReport, and IncidentReport.'),
        ('External Integration Layer', 'Google Gemini API for AI assessment, Overpass API for geospatial resource queries, and OpenStreetMap tiles for map rendering.'),
        ('Offline Layer', 'Service Worker caching strategy for UI assets (TailwindCSS, Leaflet.js, FontAwesome CDN resources), localStorage for medical profile and resource data, and IndexedDB-ready architecture for future expansion.'),
    ]
    for layer_name, layer_desc in arch_layers:
        p = doc.add_paragraph()
        run = p.add_run(f'{layer_name}: ')
        run.bold = True
        p.add_run(layer_desc)

    # -- 6. Tech Stack --
    doc.add_heading('6. Technology Stack', level=1)
    table = doc.add_table(rows=1, cols=3)
    table.style = 'Light Grid Accent 1'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Layer'
    hdr_cells[1].text = 'Technology'
    hdr_cells[2].text = 'Purpose'
    for cell in hdr_cells:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True

    tech_rows = [
        ('Backend', 'FastAPI 0.111', 'Async Python web framework'),
        ('ORM', 'SQLAlchemy 2.0', 'Database abstraction'),
        ('Database', 'SQLite', 'Lightweight, zero-config storage'),
        ('AI Engine', 'Google Gemini API', 'Severity assessment & first-aid'),
        ('Frontend', 'HTML5 + TailwindCSS', 'Responsive, mobile-first UI'),
        ('Maps', 'Leaflet.js + OSM', 'Interactive mapping'),
        ('Geospatial', 'Overpass API', 'Resource discovery from OSM'),
        ('PDF Engine', 'ReportLab', 'Incident report generation'),
        ('QR Codes', 'qrcode + Pillow', 'Medical profile QR cards'),
        ('Auth', 'JWT + bcrypt', 'Secure authentication'),
        ('PWA', 'Service Worker', 'Offline caching & install'),
        ('Deployment', 'Docker + Uvicorn', 'Containerized deployment'),
        ('Hosting', 'Render.com', 'Cloud deployment (free tier)'),
    ]
    for layer, tech, purpose in tech_rows:
        row_cells = table.add_row().cells
        row_cells[0].text = layer
        row_cells[1].text = tech
        row_cells[2].text = purpose

    # -- 7. AI Integration --
    doc.add_heading('7. AI Integration - Google Gemini API', level=1)
    doc.add_paragraph(
        'The Google Gemini API (gemini-2.0-flash model) serves as the intelligent core of ROADSoS AI, '
        'providing three critical AI capabilities:'
    )
    ai_features = [
        ('Severity Assessment', 'Analyzes incident descriptions, location context, and victim information to produce a severity rating from 1 (minor) to 5 (life-threatening). This rating determines the urgency of the response and the type of resources dispatched.'),
        ('First-Aid Guidance', 'Generates step-by-step, situation-specific first-aid instructions. These instructions are tailored to the type of injury or disaster, the available resources, and the skill level of the bystander. Instructions are delivered both as text and audio (via Web Speech API).'),
        ('Incident Summarization', 'Creates structured incident summaries for emergency services, including location details, victim count, injury descriptions, and recommended response resources.'),
    ]
    for ai_title, ai_desc in ai_features:
        p = doc.add_paragraph()
        run = p.add_run(f'{ai_title}: ')
        run.bold = True
        p.add_run(ai_desc)

    # -- 8. Offline-First PWA --
    doc.add_heading('8. Offline-First PWA Strategy', level=1)
    doc.add_paragraph(
        'Recognizing that many accident locations in the BIMSTEC region have limited connectivity, '
        'ROADSoS AI implements a comprehensive offline-first strategy:'
    )
    offline_features = [
        ('Service Worker Caching', 'All UI assets (HTML, CSS, JavaScript) and critical CDN resources (TailwindCSS, Leaflet.js, FontAwesome) are cached by the Service Worker during the first visit. Subsequent visits load entirely from cache.'),
        ('Offline Emergency Pack', 'Users can download their complete medical profile and the list of nearby emergency resources to localStorage. This data is available even without any network connection.'),
        ('PWA Installation', 'The app can be installed on the device\'s home screen via the Web App Manifest, providing a native app-like experience with a standalone display mode.'),
        ('Graceful Degradation', 'When offline, AI features gracefully degrade to show cached first-aid instructions, while the map displays cached resource locations. The SOS activation still captures GPS data and queues the request for when connectivity is restored.'),
    ]
    for of_title, of_desc in offline_features:
        p = doc.add_paragraph()
        run = p.add_run(f'{of_title}: ')
        run.bold = True
        p.add_run(of_desc)

    # -- 9. User Workflows --
    doc.add_heading('9. User Workflows', level=1)

    doc.add_heading('Emergency Workflow (Victim)', level=2)
    steps_victim = [
        'User taps the SOS button or crash is auto-detected via accelerometer',
        '5-second countdown begins with GPS lock acquisition',
        'Black Box data captured (GPS, speed, heading, network, battery)',
        'Incident created in database; AI assessment requested from Gemini',
        'Nearby hospitals, police stations, fire stations queried via Overpass API',
        'Results displayed on interactive map with distance calculations',
        'AI first-aid instructions displayed with text-to-speech audio',
        'Shareable family tracking link generated',
        'Incident report available for download (PDF/TXT/JSON)',
    ]
    for i, step in enumerate(steps_victim, 1):
        doc.add_paragraph(f'Step {i}: {step}')

    doc.add_heading('Family Tracking Workflow', level=2)
    steps_family = [
        'Victim or bystander shares the tracking link via SMS/WhatsApp',
        'Family member opens the link in any browser (no login required)',
        'Live map shows victim\'s location with a marker',
        'AI-generated response instructions are displayed',
        'Incident status updates shown in real-time',
    ]
    for i, step in enumerate(steps_family, 1):
        doc.add_paragraph(f'Step {i}: {step}')

    # -- 10. Impact & Scalability --
    doc.add_heading('10. Impact & Scalability', level=1)

    doc.add_heading('Immediate Impact', level=2)
    impacts = [
        'Reduces effective emergency response time by providing AI-guided first aid within seconds of an accident.',
        'Bridges the knowledge gap by delivering expert-level first-aid instructions to untrained bystanders.',
        'Enables offline emergency assistance in connectivity-challenged areas across the BIMSTEC region.',
        'Empowers families with real-time visibility during emergencies, reducing panic and enabling coordinated response.',
        'Provides incident heatmap data for urban planners to identify accident hotspots and improve road safety infrastructure.',
    ]
    for impact in impacts:
        doc.add_paragraph(impact, style='List Bullet')

    doc.add_heading('Scalability Path', level=2)
    scale_items = [
        'Database can migrate from SQLite to PostgreSQL for production-scale deployments.',
        'Containerized via Docker for easy deployment across cloud providers.',
        'Stateless API design enables horizontal scaling behind a load balancer.',
        'PWA architecture eliminates app store dependencies, enabling rapid distribution.',
        'Multi-language support can be added via Gemini API\'s multilingual capabilities (critical for BIMSTEC\'s diverse languages).',
    ]
    for item in scale_items:
        doc.add_paragraph(item, style='List Bullet')

    # -- 11. Future Roadmap --
    doc.add_heading('11. Future Roadmap', level=1)
    roadmap = [
        ('Phase 2 (Q3 2026)', 'Integration with national emergency services (108/112 APIs), multi-language AI responses, real-time ambulance tracking.'),
        ('Phase 3 (Q4 2026)', 'IoT integration with vehicle OBD-II ports for automatic crash data, wearable device integration for health monitoring during emergencies.'),
        ('Phase 4 (2027)', 'Cross-border emergency coordination protocol for BIMSTEC nations, integration with hospital bed availability systems, drone dispatch for remote areas.'),
    ]
    for phase, desc in roadmap:
        p = doc.add_paragraph()
        run = p.add_run(f'{phase}: ')
        run.bold = True
        p.add_run(desc)

    # -- 12. Team --
    doc.add_heading('12. Team & Acknowledgments', level=1)
    doc.add_paragraph(
        'ROADSoS AI was built with the vision of making emergency response technology accessible to '
        'every citizen across the BIMSTEC region. We acknowledge the support of Google\'s Gemini API, '
        'the OpenStreetMap community, and the Leaflet.js open-source ecosystem.'
    )
    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('GitHub: https://github.com/Harshal2007vk/ROADSOS-AI')
    run.font.color.rgb = RGBColor(59, 130, 246)

    # Save
    output_path = 'submission/ROADSoS_AI_Detailed_Document.docx'
    doc.save(output_path)
    print(f'[OK] Detailed document saved: {output_path}')
    return output_path


# ============================================================
# 2. PRESENTATION (.pptx)
# ============================================================
def create_presentation():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    BG_DARK = RGBColor(10, 10, 10)
    BG_CARD = RGBColor(30, 41, 59)
    RED = RGBColor(220, 38, 38)
    RED_LIGHT = RGBColor(239, 68, 68)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(148, 163, 184)
    GREEN = RGBColor(34, 197, 94)
    BLUE = RGBColor(59, 130, 246)
    AMBER = RGBColor(245, 158, 11)

    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def add_card(slide, left, top, width, height, title, content, title_color=RED_LIGHT):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.fill.background()
        shape.shadow.inherit = False

        # Title
        add_text_box(slide, left + 0.3, top + 0.2, width - 0.6, 0.5, title, font_size=20, color=title_color, bold=True)
        # Content
        add_text_box(slide, left + 0.3, top + 0.7, width - 0.6, height - 1.0, content, font_size=14, color=GRAY)

    def add_accent_line(slide, left, top, width):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RED
        shape.line.fill.background()

    # ========== SLIDE 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)

    # SOS circle
    sos_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(1.2), Inches(2.3), Inches(2.3))
    sos_shape.fill.solid()
    sos_shape.fill.fore_color.rgb = RED
    sos_shape.line.fill.background()
    tf = sos_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'SOS'
    p.font.size = Pt(48)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(30)

    add_text_box(slide, 1, 3.8, 11.3, 1, 'ROADSoS AI', font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.8, 11.3, 0.8, 'AI-Powered Emergency Response & Disaster Assistance Platform', font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, 5, 5.7, 3.3)
    add_text_box(slide, 1, 6.0, 11.3, 0.6, 'Road Safety Hackathon 2026  |  BIMSTEC International Track', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 6.5, 11.3, 0.6, 'github.com/Harshal2007vk/ROADSOS-AI', font_size=14, color=BLUE, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 2: Problem ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.8, 'THE PROBLEM', font_size=36, color=RED, bold=True)
    add_text_box(slide, 0.8, 1.0, 11.7, 0.6, 'Every year, 1.19 million lives are lost to road crashes. Most deaths in the BIMSTEC region are preventable.', font_size=16, color=GRAY)

    problems = [
        ('Delayed Response', '30-45 min average response time in rural areas. The Golden Hour window is missed for most victims.'),
        ('No First-Aid Knowledge', 'Bystanders lack basic first-aid skills. Improper handling worsens 40% of spinal injuries.'),
        ('Poor Connectivity', 'Highway dead zones and rural areas have no cellular coverage when accidents occur.'),
        ('Information Gaps', 'First responders arrive without knowing victim blood type, allergies, or medications.'),
        ('Resource Discovery', 'Victims cannot locate nearest hospitals in unfamiliar areas during panic.'),
        ('Family in the Dark', 'No real-time visibility for families. Hours of uncertainty and panic.'),
    ]
    positions = [(0.8, 1.8), (4.8, 1.8), (8.8, 1.8), (0.8, 4.2), (4.8, 4.2), (8.8, 4.2)]
    for (left, top), (title, desc) in zip(positions, problems):
        add_card(slide, left, top, 3.7, 2.1, title, desc)

    # ========== SLIDE 3: Solution Overview ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.8, 'OUR SOLUTION', font_size=36, color=RED, bold=True)
    add_text_box(slide, 0.8, 1.0, 11.7, 0.6, 'ROADSoS AI transforms any smartphone into a life-saving device in 3 phases.', font_size=16, color=GRAY)

    phases = [
        ('PHASE 1: DETECT', 'One-tap SOS button with 5-sec countdown\n\nAuto crash detection via accelerometer\n\nGPS lock + Black Box data capture\n\nSpeed, heading, battery, network status'),
        ('PHASE 2: ASSESS & FIND', 'Gemini AI severity assessment (1-5 scale)\n\nStep-by-step first-aid with audio\n\nNearest hospitals via Overpass API\n\nPolice & fire stations on map'),
        ('PHASE 3: COMMUNICATE', 'Shareable family tracking link\n\nNo-login public tracking page\n\nReal-time location + status updates\n\nPDF/TXT/JSON incident reports'),
    ]
    phase_colors = [RED_LIGHT, AMBER, GREEN]
    for i, ((title, desc), color) in enumerate(zip(phases, phase_colors)):
        left = 0.8 + i * 4.1
        add_card(slide, left, 1.8, 3.8, 4.8, title, desc, title_color=color)

    # ========== SLIDE 4: Key Features ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.8, 'KEY FEATURES', font_size=36, color=RED, bold=True)

    features = [
        ('One-Tap SOS', 'Big SOS button with GPS lock, countdown timer, haptic feedback'),
        ('AI First Aid', 'Gemini AI severity rating + audio-guided first-aid steps'),
        ('Crash Detection', 'Accelerometer-based auto-SOS when G-force > 4G'),
        ('Resource Map', 'Hospitals, police, fire stations on interactive Leaflet map'),
        ('Offline Pack', 'Medical profile + resources cached for zero-connectivity use'),
        ('QR Medical Card', 'Scannable QR with blood type, allergies, medications'),
        ('Family Tracking', 'Public link for real-time incident monitoring'),
        ('Black Box', 'GPS, speed, heading, network, battery recorded'),
        ('Disaster Mode', 'Flood, earthquake, cyclone - tailored AI responses'),
        ('Admin Heatmap', 'Incident density visualization for urban planning'),
        ('Report Generator', 'PDF, TXT, JSON reports for insurance & legal'),
        ('Readiness Score', '100-point gamified system for profile completion'),
    ]
    positions = []
    for row in range(3):
        for col in range(4):
            positions.append((0.5 + col * 3.2, 1.3 + row * 2.0))

    for (left, top), (title, desc) in zip(positions, features):
        add_card(slide, left, top, 2.9, 1.7, title, desc)

    # ========== SLIDE 5: Architecture ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.8, 'SYSTEM ARCHITECTURE', font_size=36, color=RED, bold=True)

    # Architecture boxes
    layers = [
        ('CLIENT (PWA)', 'HTML5 + TailwindCSS + Leaflet.js\nService Worker + localStorage\nDeviceMotionEvent (Crash Detection)\nWeb Speech API (Audio First Aid)', BLUE, 0.8, 1.5, 3.5, 2.5),
        ('FASTAPI BACKEND', 'Routes: api.py, auth.py, pages.py\nServices: user, incident, report\nModels: User, Incident, Medical\nJWT Auth + bcrypt', GREEN, 4.9, 1.5, 3.5, 2.5),
        ('EXTERNAL APIs', 'Google Gemini AI (gemini-2.0-flash)\nOverpass API (OpenStreetMap)\nOSM Tile Server (Map Rendering)', AMBER, 9.0, 1.5, 3.5, 2.5),
        ('DATA LAYER', 'SQLAlchemy 2.0 ORM\nSQLite Database\nQR Code Generator (qrcode + Pillow)\nReportLab PDF Engine', RED_LIGHT, 2.9, 4.5, 3.5, 2.5),
        ('OFFLINE LAYER', 'Service Worker Asset Caching\nlocalStorage Emergency Pack\nPWA Manifest (Installable)\nGraceful Degradation Strategy', RGBColor(168, 85, 247), 7.0, 4.5, 3.5, 2.5),
    ]
    for title, desc, color, left, top, width, height in layers:
        add_card(slide, left, top, width, height, title, desc, title_color=color)

    # ========== SLIDE 6: Tech Stack ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.8, 'TECHNOLOGY STACK', font_size=36, color=RED, bold=True)

    techs = [
        ('Backend', 'FastAPI 0.111+\nAsync Python\nPydantic Validation', BLUE),
        ('Database', 'SQLAlchemy 2.0\nSQLite\nAuto-migration', GREEN),
        ('AI Engine', 'Google Gemini\ngemini-2.0-flash\nSeverity + First Aid', RED_LIGHT),
        ('Frontend', 'HTML5 + Tailwind\nVanilla JavaScript\nJinja2 Templates', AMBER),
        ('Maps', 'Leaflet.js\nOpenStreetMap\nLeaflet.heat', RGBColor(168, 85, 247)),
        ('Geospatial', 'Overpass API\nOSM Queries\nRadius Search', BLUE),
        ('Reports', 'ReportLab PDF\nTXT + JSON\nFull Incident Log', GREEN),
        ('Auth', 'JWT (PyJWT)\nbcrypt Hashing\nSession Cookies', RED_LIGHT),
        ('PWA', 'Service Worker\nWeb Manifest\nOffline Caching', AMBER),
        ('Deploy', 'Docker\nUvicorn ASGI\nRender.com', RGBColor(168, 85, 247)),
    ]
    for i, (title, desc, color) in enumerate(techs):
        col = i % 5
        row = i // 5
        left = 0.5 + col * 2.55
        top = 1.3 + row * 3.0
        add_card(slide, left, top, 2.3, 2.7, title, desc, title_color=color)

    # ========== SLIDE 7: AI Deep Dive ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.8, 'AI INTEGRATION - GOOGLE GEMINI', font_size=36, color=RED, bold=True)
    add_text_box(slide, 0.8, 1.0, 11.7, 0.6, 'The intelligent core powering real-time emergency assessment and guidance.', font_size=16, color=GRAY)

    ai_cards = [
        ('Severity Assessment', 'Analyzes incident context and produces a 1-5 severity rating.\n\n1 = Minor scratch\n2 = Moderate injury\n3 = Serious, needs hospital\n4 = Critical, life-threatening\n5 = Mass casualty event', RED_LIGHT),
        ('First-Aid Guidance', 'Generates situation-specific instructions:\n\n- Step-by-step actions\n- Tailored to injury type\n- Audio via Web Speech API\n- Bystander skill level aware\n- Updates as situation evolves', GREEN),
        ('Incident Summary', 'Structured reports for responders:\n\n- Location + coordinates\n- Victim count + descriptions\n- Injury type classification\n- Recommended resources\n- Timeline of events', BLUE),
    ]
    for i, (title, desc, color) in enumerate(ai_cards):
        left = 0.8 + i * 4.1
        add_card(slide, left, 1.8, 3.8, 5.0, title, desc, title_color=color)

    # ========== SLIDE 8: Offline Strategy ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.8, 'OFFLINE-FIRST STRATEGY', font_size=36, color=RED, bold=True)
    add_text_box(slide, 0.8, 1.0, 11.7, 0.6, 'Designed for BIMSTEC regions where connectivity is unreliable at accident sites.', font_size=16, color=GRAY)

    offline_cards = [
        ('Service Worker', 'Caches all UI assets on first visit:\n- TailwindCSS framework\n- Leaflet.js map library\n- FontAwesome icons\n- All HTML templates\n- JavaScript modules\n\nApp loads from cache on repeat visits.', BLUE),
        ('Emergency Pack', 'User downloads to localStorage:\n- Complete medical profile\n- Blood group + allergies\n- Nearby hospitals list\n- Police station locations\n- Fire station locations\n\nAvailable with zero connectivity.', GREEN),
        ('PWA Install', 'Installable on home screen:\n- Standalone display mode\n- Custom app icon\n- Native app feel\n- No app store needed\n- Instant distribution\n\nWorks across all platforms.', AMBER),
        ('Graceful Fallback', 'When fully offline:\n- Cached first-aid shown\n- Cached resources on map\n- GPS still captures location\n- SOS queued for later send\n- Medical QR still works\n\nNever leaves user stranded.', RED_LIGHT),
    ]
    for i, (title, desc, color) in enumerate(offline_cards):
        left = 0.5 + i * 3.2
        add_card(slide, left, 1.8, 2.9, 5.0, title, desc, title_color=color)

    # ========== SLIDE 9: Impact ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.8, 'IMPACT & SCALABILITY', font_size=36, color=RED, bold=True)

    add_card(slide, 0.8, 1.3, 5.8, 5.5, 'IMMEDIATE IMPACT',
        'Reduces effective response time from 30+ min to seconds with AI first-aid\n\n'
        'Bridges first-aid knowledge gap for 2 billion+ BIMSTEC citizens\n\n'
        'Works offline in connectivity-dead zones across rural highways\n\n'
        'Empowers families with real-time tracking and peace of mind\n\n'
        'Provides incident heatmap data for urban safety planning\n\n'
        'Gamified readiness score encourages preparedness culture',
        title_color=GREEN)

    add_card(slide, 7.0, 1.3, 5.5, 5.5, 'SCALABILITY PATH',
        'SQLite -> PostgreSQL for production-scale deployments\n\n'
        'Docker containerized for any cloud provider\n\n'
        'Stateless API enables horizontal scaling\n\n'
        'PWA eliminates app store dependencies\n\n'
        'Gemini API supports BIMSTEC languages natively\n\n'
        'Open architecture for national emergency API integration',
        title_color=BLUE)

    # ========== SLIDE 10: Roadmap ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.8, 'FUTURE ROADMAP', font_size=36, color=RED, bold=True)

    roadmap = [
        ('NOW - v1.0', 'One-tap SOS + AI Assessment\nAuto crash detection\nOffline emergency pack\nFamily tracking links\nAdmin heatmap\nReport generation', GREEN),
        ('Q3 2026 - v2.0', 'National emergency API (108/112)\nMulti-language AI (BIMSTEC)\nReal-time ambulance tracking\nSMS fallback for no-internet\nCommunity responder network', AMBER),
        ('Q4 2026 - v3.0', 'Vehicle OBD-II integration\nWearable health monitoring\nHospital bed availability\nInsurance claim automation\nPredictive accident hotspots', RED_LIGHT),
        ('2027 - v4.0', 'Cross-border BIMSTEC protocol\nDrone dispatch for remote areas\nAR first-aid overlays\nBlockchain incident records\nGovernment dashboard API', RGBColor(168, 85, 247)),
    ]
    for i, (phase, desc, color) in enumerate(roadmap):
        left = 0.5 + i * 3.2
        add_card(slide, left, 1.3, 2.9, 5.5, phase, desc, title_color=color)

    # ========== SLIDE 11: Thank You ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, 0, 0, 13.333)

    # SOS circle again
    sos_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(1.0), Inches(2.3), Inches(2.3))
    sos_shape.fill.solid()
    sos_shape.fill.fore_color.rgb = RED
    sos_shape.line.fill.background()
    tf = sos_shape.text_frame
    p = tf.paragraphs[0]
    p.text = 'SOS'
    p.font.size = Pt(42)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(25)

    add_text_box(slide, 1, 3.5, 11.3, 1, 'Thank You', font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.5, 11.3, 0.6, 'Saving lives during the Golden Hour with AI', font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, 5, 5.3, 3.3)

    add_text_box(slide, 1, 5.6, 5.5, 0.5, 'GitHub', font_size=14, color=GRAY, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, 7, 5.6, 5.5, 0.5, 'github.com/Harshal2007vk/ROADSOS-AI', font_size=14, color=BLUE, alignment=PP_ALIGN.LEFT)

    add_text_box(slide, 1, 6.1, 5.5, 0.5, 'Live Demo', font_size=14, color=GRAY, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, 7, 6.1, 5.5, 0.5, 'roadsos-ai.onrender.com', font_size=14, color=BLUE, alignment=PP_ALIGN.LEFT)

    add_text_box(slide, 1, 6.6, 5.5, 0.5, 'Demo Login', font_size=14, color=GRAY, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, 7, 6.6, 5.5, 0.5, 'demo@roadsos.ai  /  demo1234', font_size=14, color=GREEN, alignment=PP_ALIGN.LEFT)

    # Save
    output_path = 'submission/ROADSoS_AI_Presentation.pptx'
    prs.save(output_path)
    print(f'[OK] Presentation saved: {output_path}')
    return output_path


# ============================================================
# 3. CODE ARCHIVE (.zip)
# ============================================================
def create_code_zip():
    output_path = 'submission/ROADSoS_AI_Code.zip'
    exclude_dirs = {'.git', '__pycache__', 'venv', '.venv', 'submission', 'node_modules'}
    exclude_files = {'.env', 'roadsos.db', 'generate_submission.py'}

    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        root = Path('.')
        for file_path in sorted(root.rglob('*')):
            # Skip excluded directories
            if any(part in exclude_dirs for part in file_path.parts):
                continue
            # Skip excluded files
            if file_path.name in exclude_files:
                continue
            # Skip directories themselves
            if file_path.is_dir():
                continue
            arcname = str(Path('ROADSoS_AI') / file_path)
            zipf.write(file_path, arcname)
            print(f'  Added: {arcname}')

    print(f'[OK] Code archive saved: {output_path}')
    return output_path


# ============================================================
# MAIN
# ============================================================
if __name__ == '__main__':
    os.makedirs('submission', exist_ok=True)
    print('='*60)
    print('  ROADSoS AI - Hackathon Submission Generator')
    print('='*60)
    print()

    print('[1/3] Generating detailed document...')
    create_detailed_document()
    print()

    print('[2/3] Generating presentation...')
    create_presentation()
    print()

    print('[3/3] Creating code archive...')
    create_code_zip()
    print()

    print('='*60)
    print('  ALL SUBMISSION FILES READY!')
    print('='*60)
    print()
    print('Files in submission/ folder:')
    for f in sorted(Path('submission').iterdir()):
        size_mb = f.stat().st_size / (1024 * 1024)
        print(f'  {f.name:45s} ({size_mb:.2f} MB)')
    print()
    print('Upload these to Unstop:')
    print('  1. Detailed Document -> ROADSoS_AI_Detailed_Document.docx')
    print('  2. Presentation PPT  -> ROADSoS_AI_Presentation.pptx')
    print('  3. Code              -> ROADSoS_AI_Code.zip')
